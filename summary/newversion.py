# =======================================
# newversion.py – Extract lecture text from uploaded file OR text input
# Supports PDF, TXT, DOCX, PPTX/PPT, and direct text input
# =======================================

import os
import sys
import re
from openai import OpenAI
import nltk
from nltk.corpus import stopwords

# Optional: extra libraries for file handling
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Download stopwords if not already
nltk.download("stopwords")
stop_words = set(stopwords.words("english"))

# -------------------------
# GitHub AI Setup
# -------------------------
os.environ["GITHUB_TOKEN"]= "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"
token = os.environ["GITHUB_TOKEN"]
endpoint = "https://models.github.ai/inference"
client = OpenAI(api_key=token, base_url=endpoint)

# -------------------------
# Read file path or text flag from argument
# -------------------------
if len(sys.argv) < 2:
    raise ValueError("Usage: python newversion.py <lecture_file_path> OR python newversion.py --text <input_text>")

# Check if it's text input mode
if sys.argv[1] == "--text":
    if len(sys.argv) < 3:
        raise ValueError("Text input mode requires text content")
    # Text mode: read from argument
    raw_text = sys.argv[2]
    output_dir = os.getcwd()
else:
    # File mode: read from file
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"{file_path} not found")
    output_dir = os.path.dirname(file_path)

# -------------------------
# Extract text functions
# -------------------------
def extract_pdf_text(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_txt_text(txt_path):
    with open(txt_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_docx_text(docx_path):
    doc = Document(docx_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx_text(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path):
    ext = file_path.split(".")[-1].lower()
    if ext == "pdf":
        return extract_pdf_text(file_path)
    elif ext == "txt":
        return extract_txt_text(file_path)
    elif ext == "docx":
        return extract_docx_text(file_path)
    elif ext in ["ppt", "pptx"]:
        return extract_pptx_text(file_path)
    else:
        raise ValueError("Unsupported file type")

# -------------------------
# Preprocess text
# -------------------------
def preprocess_text(text):
    text = text.replace("�", " ")
    text = text.encode("ascii", errors="ignore").decode()
    text = re.sub(r'\s+', ' ', text).strip()

    patterns_to_remove = [
        r'\bChapter\s+\d+\b',
        r'\bIntroduction\b',
        r'\bConclusion\b',
        r'\bSummary\b',
        r'\bReferences\b',
        r'\bBibliography\b',
        r'\bTable of Contents\b',
        r'\bAppendix\b',
        r'\bFigure\s+\d+\b',
        r'\bTable\s+\d+\b',
    ]
    for pattern in patterns_to_remove:
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)

    words = text.split()
    words = [w for w in words if w.lower() not in stop_words]
    return " ".join(words)

# -------------------------
# Clean and split sentences using GitHub AI
# -------------------------
def clean_and_split(text, model_name="openai/gpt-4o"):
    prompt = f"""
You are a text cleaning assistant.
Clean the text below by:
- Removing headings, numbers, references, and URLs
- Fixing broken sentences
- Splitting into proper complete sentences
- Return only the cleaned sentences as a list

Text:
{text}
"""
    response = client.chat.completions.create(
        model=model_name,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

# -------------------------
# Main processing
# -------------------------
if sys.argv[1] != "--text":
    raw_text = extract_text(file_path)

preprocessed_text = preprocess_text(raw_text)
cleaned_sentences = clean_and_split(preprocessed_text)

# -------------------------
# Save output
# -------------------------
output_file = os.path.join(output_dir, "ai_extracted.txt")
with open(output_file, "w", encoding="utf-8") as f:
    f.write(cleaned_sentences)

print(f"✅ ai_extracted.txt saved at {output_file}")